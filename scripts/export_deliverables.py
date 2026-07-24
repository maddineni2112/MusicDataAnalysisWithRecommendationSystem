from __future__ import annotations

import re
import textwrap
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import ListFlowable, ListItem, Paragraph, SimpleDocTemplate, Spacer

ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
EXPORTS = DOCS / "exports"


def clean_inline(text: str) -> str:
    text = re.sub(r"`([^`]+)`", r"\1", text)
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    text = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
    return text


def markdown_to_pdf(source: Path, target: Path) -> None:
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name="ReportTitle", parent=styles["Title"], fontSize=22, leading=28, spaceAfter=18))
    styles.add(ParagraphStyle(name="Heading2Custom", parent=styles["Heading2"], fontSize=15, leading=18, spaceBefore=14, spaceAfter=8))
    styles.add(ParagraphStyle(name="BodyCustom", parent=styles["BodyText"], fontSize=10.5, leading=14, spaceAfter=7))
    styles.add(ParagraphStyle(name="BulletCustom", parent=styles["BodyText"], fontSize=10, leading=13))
    doc = SimpleDocTemplate(str(target), pagesize=letter, rightMargin=0.7 * inch, leftMargin=0.7 * inch, topMargin=0.7 * inch, bottomMargin=0.7 * inch)
    story = []
    bullet_buffer = []

    def flush_bullets() -> None:
        nonlocal bullet_buffer
        if bullet_buffer:
            story.append(ListFlowable([ListItem(Paragraph(clean_inline(item), styles["BulletCustom"])) for item in bullet_buffer], bulletType="bullet"))
            story.append(Spacer(1, 0.08 * inch))
            bullet_buffer = []

    for raw_line in source.read_text(encoding="utf-8").splitlines():
        line = raw_line.strip()
        if not line:
            flush_bullets()
            continue
        if line.startswith("```"):
            flush_bullets()
            continue
        if line.startswith("# "):
            flush_bullets()
            story.append(Paragraph(clean_inline(line[2:]), styles["ReportTitle"]))
        elif line.startswith("## "):
            flush_bullets()
            story.append(Paragraph(clean_inline(line[3:]), styles["Heading2Custom"]))
        elif line.startswith("- "):
            bullet_buffer.append(line[2:])
        else:
            flush_bullets()
            story.append(Paragraph(clean_inline(line), styles["BodyCustom"]))
    flush_bullets()
    doc.build(story)


def parse_deck(source: Path) -> list[dict[str, list[str] | str]]:
    slides = []
    current: dict[str, list[str] | str] | None = None
    paragraph_buffer: list[str] = []

    def flush_paragraph() -> None:
        nonlocal paragraph_buffer
        if current is not None and paragraph_buffer:
            current["body"].append(" ".join(paragraph_buffer))
            paragraph_buffer = []

    for raw_line in source.read_text(encoding="utf-8").splitlines():
        line = raw_line.strip()
        if line.startswith("## Slide "):
            flush_paragraph()
            title = line.split(":", 1)[1].strip() if ":" in line else line.replace("##", "").strip()
            current = {"title": title, "body": [], "bullets": []}
            slides.append(current)
        elif current is None or line.startswith("# ") or not line:
            flush_paragraph()
        elif line.startswith("- "):
            flush_paragraph()
            current["bullets"].append(line[2:])
        elif re.match(r"^\d+\.\s+", line):
            flush_paragraph()
            current["bullets"].append(re.sub(r"^\d+\.\s+", "", line))
        else:
            paragraph_buffer.append(line)
    flush_paragraph()
    return slides


def add_wrapped_lines(text_frame, lines: list[str], font_size: int = 20) -> None:
    first = True
    for line in lines:
        wrapped = textwrap.wrap(clean_inline(line), width=72) or [""]
        for index, part in enumerate(wrapped):
            paragraph = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
            paragraph.text = part
            paragraph.font.size = Pt(font_size)
            paragraph.font.color.rgb = RGBColor.from_string("1f2937")
            paragraph.level = 0
            first = False
            if index > 0:
                paragraph.level = 1


def markdown_to_pptx(source: Path, target: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_data in parse_deck(source):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor.from_string("f8fafc")

        title_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.45), Inches(12), Inches(0.75))
        title_frame = title_box.text_frame
        title_frame.text = clean_inline(str(slide_data["title"]))
        title_frame.paragraphs[0].font.size = Pt(30)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor.from_string("111827")

        body_box = slide.shapes.add_textbox(Inches(0.85), Inches(1.5), Inches(11.7), Inches(5.3))
        body_frame = body_box.text_frame
        body_frame.word_wrap = True
        lines = [str(item) for item in slide_data["body"]]
        bullets = [str(item) for item in slide_data["bullets"]]
        add_wrapped_lines(body_frame, lines, 20)
        for bullet in bullets:
            paragraph = body_frame.add_paragraph()
            paragraph.text = clean_inline(bullet)
            paragraph.font.size = Pt(19)
            paragraph.font.color.rgb = RGBColor.from_string("1f2937")
            paragraph.level = 0

        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.18), Inches(13.333), Inches(0.08))
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor.from_string("1db954")
        accent.line.color.rgb = RGBColor.from_string("1db954")

    prs.save(target)


def main() -> None:
    EXPORTS.mkdir(parents=True, exist_ok=True)
    markdown_to_pdf(DOCS / "technical-report-v2.md", EXPORTS / "Indian_Music_Intelligence_Platform_v2_Report.pdf")
    markdown_to_pptx(DOCS / "pitch-deck-v2.md", EXPORTS / "Indian_Music_Intelligence_Platform_v2_Pitch_Deck.pptx")
    print(f"Exported deliverables to {EXPORTS}")


if __name__ == "__main__":
    main()
